import csv
import io
import json
import tempfile
from pathlib import Path
from unittest.mock import patch

from django.test import SimpleTestCase, override_settings
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from apps.collab import xiaoce_file_artifacts
from apps.collab.xiaoce_file_artifacts import (
    FORMAT_SPECS,
    ArtifactRequest,
    detect_file_artifact_requests,
    extract_explicit_file_content,
    maybe_generate_file_artifacts,
)


class XiaoceFileArtifactIntentTests(SimpleTestCase):
    def test_detects_multiple_named_formats_and_safe_shared_filename(self):
        requests = detect_file_artifact_requests(
            "请生成 md、PDF、Word 和 PPT 文件，内容写 hello，"
            '文件名为 "../../季度总结.md"',
        )

        self.assertEqual(
            requests,
            [
                ArtifactRequest("md", "季度总结.md"),
                ArtifactRequest("pdf", "季度总结.pdf"),
                ArtifactRequest("docx", "季度总结.docx"),
                ArtifactRequest("pptx", "季度总结.pptx"),
            ],
        )

    def test_explicit_extension_in_filename_can_identify_the_format(self):
        requests = detect_file_artifact_requests(
            "Generate a file, file name: report.json, content: hello",
        )
        self.assertEqual(requests, [ArtifactRequest("json", "report.json")])

    def test_code_format_aliases_are_supported(self):
        cases = {
            "生成 Python 脚本": "py",
            "创建一个 JavaScript 文件": "js",
            "生成 TypeScript 文件": "ts",
            "创建 JSX 文件": "jsx",
            "创建 TSX 文件": "tsx",
            "制作 CSS 样式表文件": "css",
            "导出 SQL 脚本": "sql",
            "生成 YAML 配置文件": "yaml",
            "生成 yml 文件": "yaml",
            "生成 XML 文档": "xml",
        }
        for message, expected in cases.items():
            with self.subTest(message=message):
                requests = detect_file_artifact_requests(message)
                self.assertEqual([item.format for item in requests], [expected])

    def test_weak_production_verbs_require_an_explicit_file_or_count_target(self):
        self.assertEqual(
            detect_file_artifact_requests("帮我做一个 PDF，内容写 hello"),
            [ArtifactRequest("pdf", "xiaoce-export.pdf")],
        )
        self.assertEqual(
            detect_file_artifact_requests("Please make a PDF with content: hello"),
            [ArtifactRequest("pdf", "xiaoce-export.pdf")],
        )
        self.assertEqual(
            detect_file_artifact_requests("输出 JSON"),
            [],
        )

    def test_read_resend_inability_and_response_format_requests_do_not_generate(self):
        cases = (
            "读取这个 PDF 并分析",
            "把刚才生成的 PDF 文件发给我",
            "返回之前创建的 Word 文档",
            "为什么不能生成 PDF？",
            "小策无法创建 docx 文件吗？",
            "不要生成 HTML 文件",
            "Please do not create a markdown file.",
            "Why can't you generate a PDF?",
            "输出 JSON，不要 markdown。",
            "Output JSON only.",
        )
        for message in cases:
            with self.subTest(message=message):
                self.assertEqual(detect_file_artifact_requests(message), [])

    def test_negated_first_clause_does_not_hide_a_later_positive_request(self):
        requests = detect_file_artifact_requests(
            "不要生成 PDF，改为生成 Markdown 文件，内容写 hello",
        )
        self.assertEqual(requests, [ArtifactRequest("md", "xiaoce-export.md")])

    def test_input_formats_are_not_mistaken_for_requested_outputs(self):
        cases = (
            (
                "请读取附件中的 Excel 数据，并生成一份 PDF 分析报告",
                ["pdf"],
            ),
            (
                "请生成一个 Word 文档，内容为对上传 PDF 的总结",
                ["docx"],
            ),
            (
                "请阅读下面的 Python 代码并生成 Word 分析报告\n"
                "```python\nprint('hello')\n```",
                ["docx"],
            ),
            (
                "请生成一份基于 Excel 数据的 PDF 报告",
                ["pdf"],
            ),
        )
        for message, expected in cases:
            with self.subTest(message=message):
                self.assertEqual(
                    [item.format for item in detect_file_artifact_requests(message)],
                    expected,
                )

    def test_derived_content_instructions_do_not_skip_model_reasoning(self):
        self.assertIsNone(
            extract_explicit_file_content(
                "请生成一个 Word 文档，内容为对上传 PDF 的总结",
            ),
        )
        self.assertIsNone(
            extract_explicit_file_content(
                "请阅读下面的 Python 代码并生成 Word 分析报告\n"
                "```python\nprint('hello')\n```",
            ),
        )
        self.assertEqual(
            extract_explicit_file_content(
                "请生成 Word 文档，内容为“对上传 PDF 的总结”",
            ),
            "对上传 PDF 的总结",
        )

    def test_more_than_five_formats_is_rejected(self):
        with self.assertRaisesRegex(ValueError, "最多生成 5 个"):
            detect_file_artifact_requests(
                "生成 md、txt、html、json、csv 和 pdf 文件",
            )

    def test_extracts_plain_and_fenced_literal_content(self):
        self.assertEqual(
            extract_explicit_file_content(
                "生成 Markdown 文件，内容里面写一个 hello 就好，文件名为 a.md",
            ),
            "hello",
        )
        self.assertEqual(
            extract_explicit_file_content(
                "生成 Python 文件：\n```python\nprint('hello')\n```",
            ),
            "print('hello')",
        )
        self.assertEqual(
            extract_explicit_file_content(
                "请在上传的 Excel 的 B2 单元格写入 hello，"
                "然后把文件发给我，并导出 Excel",
            ),
            "hello",
        )


class XiaoceFileArtifactGenerationTests(SimpleTestCase):
    def generate(
        self,
        root: str,
        message: str,
        *,
        model_reply: str | None = None,
        user_id: int = 17,
    ) -> list[dict]:
        with override_settings(CHAT_ATTACHMENTS_ROOT=Path(root)):
            return maybe_generate_file_artifacts(
                user_id=user_id,
                request_text=message,
                model_reply=model_reply,
            )

    def read_artifact(self, root: str, item: dict, user_id: int = 17) -> bytes:
        path = Path(root) / str(user_id) / item["id"]
        self.assertTrue(path.is_file())
        self.assertTrue(path.resolve().is_relative_to((Path(root) / str(user_id)).resolve()))
        return path.read_bytes()

    def test_generates_five_text_data_formats_with_valid_content_and_mime(self):
        with tempfile.TemporaryDirectory() as tmp:
            generated = self.generate(
                tmp,
                "生成 md、txt、html、json 和 csv 文件，内容写 hello",
            )
            by_format = {
                Path(item["name"]).suffix.lstrip("."): item
                for item in generated
            }

            self.assertEqual(list(by_format), ["md", "txt", "html", "json", "csv"])
            for file_format, item in by_format.items():
                self.assertEqual(item["mime"], FORMAT_SPECS[file_format].mime)
                self.assertEqual(item["size"], len(self.read_artifact(tmp, item)))
                self.assertTrue(item["is_file"])
                self.assertFalse(item["is_image"])

            self.assertEqual(self.read_artifact(tmp, by_format["md"]), b"hello\n")
            self.assertEqual(self.read_artifact(tmp, by_format["txt"]), b"hello\n")
            html_text = self.read_artifact(tmp, by_format["html"]).decode("utf-8")
            self.assertIn("<!doctype html>", html_text)
            self.assertIn("<pre>hello</pre>", html_text)
            json_value = json.loads(
                self.read_artifact(tmp, by_format["json"]).decode("utf-8"),
            )
            self.assertEqual(json_value, {"content": "hello"})
            csv_text = self.read_artifact(tmp, by_format["csv"]).decode("utf-8-sig")
            self.assertEqual(list(csv.reader(io.StringIO(csv_text))), [["hello"]])

    def test_model_reply_is_used_when_the_request_has_no_literal_content(self):
        with tempfile.TemporaryDirectory() as tmp:
            generated = self.generate(
                tmp,
                "请创建一份 Markdown 报告文件",
                model_reply="# 周报\n\n本周完成 3 项。",
            )
            content = self.read_artifact(tmp, generated[0]).decode("utf-8")

        self.assertIn("# 周报", content)
        self.assertIn("本周完成 3 项", content)

    def test_missing_literal_and_model_content_produces_no_fake_file(self):
        with tempfile.TemporaryDirectory() as tmp:
            generated = self.generate(tmp, "请创建一份 Markdown 报告文件")
            self.assertEqual(generated, [])
            self.assertFalse((Path(tmp) / "17").exists())

    def test_fenced_code_is_extracted_and_filename_cannot_escape_user_root(self):
        with tempfile.TemporaryDirectory() as tmp:
            generated = self.generate(
                tmp,
                "创建 Python 脚本，文件名为 \"../../evil.py\"\n"
                "```python\nprint('hello')\n```",
            )
            item = generated[0]
            content = self.read_artifact(tmp, item).decode("utf-8")

        self.assertEqual(item["name"], "evil.py")
        self.assertNotIn("..", item["id"])
        self.assertEqual(content, "print('hello')\n")

    def test_all_safe_code_formats_are_plain_non_executable_attachments(self):
        cases = {
            "js": ("生成 JavaScript 文件，内容写 console.log('hello')", "console.log"),
            "ts": ("生成 TypeScript 文件，内容写 const value: string = 'hello'", "const value"),
            "jsx": ("生成 JSX 文件，内容写 <div>hello</div>", "<div>"),
            "tsx": ("生成 TSX 文件，内容写 <div>hello</div>", "<div>"),
            "css": ("生成 CSS 文件，内容写 body { color: red; }", "color"),
            "sql": ("导出 SQL 脚本，内容写 SELECT 'hello';", "SELECT"),
            "yaml": ("生成 YAML 配置文件，内容写 greeting: hello", "greeting"),
            "xml": ("生成 XML 文档，内容写 hello", "<content>hello</content>"),
        }
        for file_format, (message, needle) in cases.items():
            with self.subTest(file_format=file_format), tempfile.TemporaryDirectory() as tmp:
                item = self.generate(tmp, message)[0]
                content = self.read_artifact(tmp, item).decode("utf-8")
                self.assertEqual(Path(item["name"]).suffix, f".{file_format}")
                self.assertIn(needle, content)

    def test_json_fence_is_parsed_and_pretty_printed(self):
        with tempfile.TemporaryDirectory() as tmp:
            item = self.generate(
                tmp,
                "生成 JSON 文件\n```json\n{\"hello\": [1, 2]}\n```",
            )[0]
            payload = json.loads(self.read_artifact(tmp, item).decode("utf-8"))
        self.assertEqual(payload, {"hello": [1, 2]})

    def test_csv_cells_are_neutralized_against_spreadsheet_formulas(self):
        with tempfile.TemporaryDirectory() as tmp:
            item = self.generate(
                tmp,
                "生成 CSV 文件\n```csv\nname,value\none,=2+2\ntwo,+cmd\n"
                "three,-10\nfour,@SUM(A1:A2)\n```",
            )[0]
            rows = list(
                csv.reader(
                    io.StringIO(
                        self.read_artifact(tmp, item).decode("utf-8-sig"),
                    ),
                ),
            )

        self.assertEqual(rows[1][1], "'=2+2")
        self.assertEqual(rows[2][1], "'+cmd")
        self.assertEqual(rows[3][1], "'-10")
        self.assertEqual(rows[4][1], "'@SUM(A1:A2)")

    def test_html_is_static_and_removes_script_event_and_javascript_url(self):
        with tempfile.TemporaryDirectory() as tmp:
            item = self.generate(
                tmp,
                "生成 HTML 文件\n```html\n"
                '<html><head><meta http-equiv="refresh" content="0;url=https://evil">'
                '<style>@import "https://evil/style.css";</style></head>'
                '<body onload="steal()"><script>alert(1)</script>'
                '<a href=javascript:alert(1)>one</a>'
                '<a href="java&#x73;cript:steal()">hello</a>'
                '<img src="https://evil/pixel.png"></body></html>\n```',
            )[0]
            content = self.read_artifact(tmp, item).decode("utf-8")

        self.assertNotIn("<script", content.lower())
        self.assertNotIn("onload", content.lower())
        self.assertNotIn("javascript:", content.lower())
        self.assertNotIn("http-equiv=\"refresh\"", content.lower())
        self.assertNotIn("evil/style.css", content.lower())
        self.assertNotIn("evil/pixel.png", content.lower())
        self.assertIn("Content-Security-Policy", content)
        self.assertIn(">hello</a>", content)

    def test_docx_pdf_and_pptx_are_openable_and_contain_the_requested_text(self):
        with tempfile.TemporaryDirectory() as tmp:
            generated = self.generate(
                tmp,
                "生成 Word、PDF 和 PPT 文件，内容写 你好 hello",
            )
            by_suffix = {Path(item["name"]).suffix: item for item in generated}

            docx_path = Path(tmp) / "17" / by_suffix[".docx"]["id"]
            self.assertTrue(docx_path.read_bytes().startswith(b"PK"))
            document = Document(docx_path)
            self.assertIn("hello", "\n".join(p.text for p in document.paragraphs))

            pdf_data = self.read_artifact(tmp, by_suffix[".pdf"])
            self.assertTrue(pdf_data.startswith(b"%PDF"))
            pdf = PdfReader(io.BytesIO(pdf_data))
            self.assertEqual(len(pdf.pages), 1)
            self.assertIn("你好", pdf.pages[0].extract_text())
            self.assertIn("hello", pdf.pages[0].extract_text())

            pptx_path = Path(tmp) / "17" / by_suffix[".pptx"]["id"]
            self.assertTrue(pptx_path.read_bytes().startswith(b"PK"))
            presentation = Presentation(pptx_path)
            slide_text = "\n".join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
            self.assertIn("hello", slide_text)

    def test_xlsx_is_delegated_and_custom_filename_is_applied_safely(self):
        with tempfile.TemporaryDirectory() as tmp:
            generated = self.generate(
                tmp,
                "生成 Excel 文件，内容里面写一个 hello，"
                '文件名为 "../../greeting.xlsx"',
            )
            item = generated[0]
            path = Path(tmp) / "17" / item["id"]
            workbook = load_workbook(path, data_only=False)
            try:
                value = workbook.active["A1"].value
            finally:
                workbook.close()

        self.assertEqual(item["name"], "greeting.xlsx")
        self.assertEqual(value, "hello")

    def test_renderer_exception_is_not_swallowed_and_partial_files_are_cleaned(self):
        with tempfile.TemporaryDirectory() as tmp:
            with (
                override_settings(CHAT_ATTACHMENTS_ROOT=Path(tmp)),
                patch.dict(
                    xiaoce_file_artifacts._RENDERERS,
                    {"pdf": lambda _content: (_ for _ in ()).throw(RuntimeError("boom"))},
                ),
            ):
                with self.assertRaisesRegex(RuntimeError, "boom"):
                    maybe_generate_file_artifacts(
                        user_id=17,
                        request_text="生成 Markdown 和 PDF 文件，内容写 hello",
                        model_reply=None,
                    )
                remaining = list((Path(tmp) / "17").glob("*"))

        self.assertEqual(remaining, [])
