"""Deterministic, local file artifacts for Xiaoce chat requests.

The language model supplies *content*, never a filesystem path.  This module
recognises only explicit file-production requests, renders a small allow-list
of formats, and atomically stores the result in the current user's attachment
directory.  Returned dictionaries intentionally match the ``generated_files``
shape consumed by the collaboration worker.
"""
from __future__ import annotations

import csv
import html
import io
import json
import os
import re
import tempfile
import unicodedata
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Callable

from apps.core.attachments import attachments_root, resolve_attachment_path

from .xiaoce_artifacts import (
    extract_excel_a1_content,
    maybe_generate_excel_artifact,
)


MAX_FILE_ARTIFACTS = 5
MAX_GENERATED_FILE_BYTES = 20 * 1024 * 1024
MAX_ARTIFACT_TEXT_CHARS = 1_000_000
DEFAULT_FILE_STEM = "xiaoce-export"


@dataclass(frozen=True)
class ArtifactFormat:
    extension: str
    mime: str
    aliases: tuple[str, ...]
    text_file: bool = False


@dataclass(frozen=True)
class ArtifactRequest:
    """One validated output request detected from the user's message."""

    format: str
    filename: str


FORMAT_SPECS: dict[str, ArtifactFormat] = {
    "md": ArtifactFormat(
        "md",
        "text/markdown; charset=utf-8",
        (r"md", r"markdown", r"markdown\s*(?:文件|文档)", r"马克当"),
        True,
    ),
    "txt": ArtifactFormat(
        "txt",
        "text/plain; charset=utf-8",
        (r"txt", r"纯文本(?:文件|文档)?", r"文本文件", r"plain\s*text"),
        True,
    ),
    "html": ArtifactFormat(
        "html",
        "text/html; charset=utf-8",
        (r"html?", r"网页(?:文件|页面)?", r"web\s*page"),
        True,
    ),
    "json": ArtifactFormat(
        "json",
        "application/json; charset=utf-8",
        (r"json", r"json\s*(?:文件|数据)"),
        True,
    ),
    "csv": ArtifactFormat(
        "csv",
        "text/csv; charset=utf-8",
        (r"csv", r"逗号分隔(?:文件|表格)?", r"comma[-\s]*separated"),
        True,
    ),
    "docx": ArtifactFormat(
        "docx",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        (r"docx?", r"word(?:\s*(?:文件|文档|document))?", r"微软文档"),
    ),
    "pdf": ArtifactFormat(
        "pdf",
        "application/pdf",
        (r"pdf", r"便携式文档"),
    ),
    "pptx": ArtifactFormat(
        "pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        (
            r"pptx?",
            r"power\s*point",
            r"powerpoint",
            r"幻灯片",
            r"演示文稿",
            r"演示文件",
        ),
    ),
    "xlsx": ArtifactFormat(
        "xlsx",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        (r"xlsx?", r"excel", r"电子表格", r"工作簿", r"spreadsheet", r"workbook"),
    ),
    "py": ArtifactFormat(
        "py",
        "text/x-python; charset=utf-8",
        (r"py", r"python(?:\s*(?:文件|脚本|代码))?"),
        True,
    ),
    "js": ArtifactFormat(
        "js",
        "text/javascript; charset=utf-8",
        (r"js", r"javascript(?:\s*(?:文件|脚本|代码))?"),
        True,
    ),
    "ts": ArtifactFormat(
        "ts",
        "text/typescript; charset=utf-8",
        (r"ts", r"typescript(?:\s*(?:文件|脚本|代码))?"),
        True,
    ),
    "jsx": ArtifactFormat(
        "jsx",
        "text/jsx; charset=utf-8",
        (r"jsx", r"react\s+jsx"),
        True,
    ),
    "tsx": ArtifactFormat(
        "tsx",
        "text/tsx; charset=utf-8",
        (r"tsx", r"react\s+tsx"),
        True,
    ),
    "css": ArtifactFormat(
        "css",
        "text/css; charset=utf-8",
        (r"css", r"样式表(?:文件)?", r"stylesheet"),
        True,
    ),
    "sql": ArtifactFormat(
        "sql",
        "text/x-sql; charset=utf-8",
        (r"sql", r"sql\s*(?:文件|脚本|代码)"),
        True,
    ),
    "yaml": ArtifactFormat(
        "yaml",
        "application/yaml; charset=utf-8",
        (r"yaml", r"yml", r"yaml\s*(?:文件|配置)"),
        True,
    ),
    "xml": ArtifactFormat(
        "xml",
        "application/xml; charset=utf-8",
        (r"xml", r"xml\s*(?:文件|文档)"),
        True,
    ),
}

_FORMAT_PATTERNS: dict[str, re.Pattern[str]] = {
    key: re.compile(
        "|".join(
            rf"(?<![A-Za-z0-9_])(?:{alias})(?![A-Za-z0-9_])"
            for alias in spec.aliases
        ),
        re.IGNORECASE,
    )
    for key, spec in FORMAT_SPECS.items()
}

_STRONG_PRODUCTION = (
    r"(?:重新\s*)?(?:生成|产出|创建|制作|导出)(?!\s*(?:的|过的|出来的|好的))"
    r"|(?:保存为|另存为|转换为|转成)"
    r"|\b(?:generate|produce|create|export|save\s+as|convert\s+to)\b"
)
_WEAK_PRODUCTION = (
    r"(?:写|做|输出)(?!\s*(?:的|过的|出来的|好的))"
    r"|\b(?:write|make|output)\b"
)
_ACTION_RE = re.compile(
    rf"(?P<strong>{_STRONG_PRODUCTION})|(?P<weak>{_WEAK_PRODUCTION})",
    re.IGNORECASE,
)
_FILE_NOUN_RE = re.compile(
    r"(?:文件|附件|文档|报告|网页|页面|脚本|代码|表格|工作簿|幻灯片|演示文稿)"
    r"|\b(?:file|attachment|document|report|page|script|spreadsheet|"
    r"workbook|presentation|deck)\b",
    re.IGNORECASE,
)
_COUNT_TARGET_RE = re.compile(
    r"(?:一|1)\s*(?:个|份|篇|张|套)"
    r"|\b(?:a|an|one)\b",
    re.IGNORECASE,
)
_NEGATED_PRODUCTION_RE = re.compile(
    r"(?:"
    r"(?:不要|别|无需|不用|不需要|禁止|停止|取消).{0,18}"
    r"(?:生成|产出|创建|制作|导出|保存|转换|写|做|输出)"
    r"|(?:为什么|为何|怎么会).{0,20}(?:(?<!能)不能|无法|没法|不支持|做不到)"
    r"|(?:无法|没法|不支持|做不到).{0,18}"
    r"(?:生成|产出|创建|制作|导出|保存|转换|写|做|输出)"
    r"|\b(?:do\s+not|don't|dont|never|no\s+need\s+to)\b.{0,20}"
    r"\b(?:generate|produce|create|export|save|convert|write|make|output)\b"
    r"|\bwhy\b.{0,24}\b(?:cannot|can't|cant|unable)\b"
    r")",
    re.IGNORECASE | re.DOTALL,
)
_SENTENCE_BOUNDARY_RE = re.compile(r"[。！？!?；;\n]+")
_COMMA_RE = re.compile(r"[,，]")
_INPUT_CONTEXT_RE = re.compile(
    r"(?:读取|阅读|查看|解析|分析|总结|提取|基于|根据|参考|上传|附件)"
    r"|\b(?:read|inspect|parse|analy[sz]e|summari[sz]e|extract|"
    r"based\s+on|according\s+to|attached|uploaded)\b",
    re.IGNORECASE,
)
_OUTPUT_SCOPE_BOUNDARY_RE = re.compile(
    r"[,，]\s*(?:"
    r"(?:内容|正文|数据|代码)(?:里面|里|中)?\s*"
    r"(?:写(?:入|上)?|填写|输入|为|是|包含|放入?|如下|：|:)"
    r"|(?:根据|基于|参考|读取|阅读|查看|解析|分析|总结|提取|使用)"
    r"|\b(?:with\s+(?:the\s+)?(?:content|text|body|data|code)|"
    r"based\s+on|using|from\s+(?:the\s+)?(?:attached|uploaded))\b"
    r")",
    re.IGNORECASE,
)
_FORMAT_INPUT_REGION_RE = re.compile(
    r"(?:基于|根据|参考|来自).{0,64}?的"
    r"|\b(?:based\s+on|according\s+to|using\s+(?:the\s+)?"
    r"(?:attached|uploaded))\b.*$",
    re.IGNORECASE,
)

_FENCED_BLOCK_RE = re.compile(
    r"```(?P<language>[A-Za-z0-9_+#.-]*)[ \t]*\r?\n"
    r"(?P<body>.*?)"
    r"\r?\n?```",
    re.DOTALL,
)
_CONTENT_PATTERNS = (
    re.compile(
        r"(?:内容|正文|数据|代码)(?:里面|里|中)?\s*"
        r"(?:写(?:入|上)?|填写|输入|为|是|包含|放入?|：|:)\s*"
        r"(?:一个|一段|一条)?\s*(?P<value>.+)",
        re.IGNORECASE | re.DOTALL,
    ),
    re.compile(
        r"(?:with\s+)?(?:the\s+)?(?:content|text|body|data|code)\s*"
        r"(?:should\s*)?(?:be|say|contain|read|write|:|=)\s*"
        r"(?P<value>.+)",
        re.IGNORECASE | re.DOTALL,
    ),
)
_FOLLOW_UP_RE = re.compile(
    r"(?:"
    r"\s*(?:[，,；;。]\s*)?(?:然后|随后|接着)"
    r"|\s*(?:[，,；;。]\s*)?并(?:且)?\s*(?:把|将)?\s*(?:生成的|这个|该)?\s*"
    r"文件\s*(?:发|发送|给|返回|回传|保存|下载|命名)"
    r"|\s*(?:[，,；;。]\s*)?(?:文件名|档案名)\s*(?:为|叫|是|：|:)"
    r"|\s*(?:,\s*)?\bthen\b"
    r"|\s*(?:,\s*)?\band\s+(?:send|return|give|name|save|download)\b"
    r"|\s*(?:,\s*)?\bfile\s*name\s*(?:is|should\s+be|:)"
    r")",
    re.IGNORECASE | re.DOTALL,
)
_FILENAME_RE = re.compile(
    r"(?:"
    r"(?:文件名|档案名)\s*(?:为|叫|是|：|:)\s*"
    r"|(?:file\s*name|filename)\s*(?:is|should\s+be|:|=)\s*"
    r"|(?:named|name\s+it)\s+"
    r")"
    r"(?:[\"“'‘](?P<quoted>[^\"”'’\r\n]{1,180})[\"”'’]"
    r"|(?P<bare>[^\s,，;；。!?！？]{1,180}))",
    re.IGNORECASE,
)

_FENCE_LANGUAGE_FORMATS: dict[str, set[str]] = {
    "md": {"md", "markdown"},
    "txt": {"txt", "text", "plaintext"},
    "html": {"html", "htm"},
    "json": {"json", "jsonc"},
    "csv": {"csv"},
    "py": {"py", "python"},
    "js": {"js", "javascript", "node"},
    "ts": {"ts", "typescript"},
    "jsx": {"jsx"},
    "tsx": {"tsx"},
    "css": {"css"},
    "sql": {"sql"},
    "yaml": {"yaml", "yml"},
    "xml": {"xml"},
}


def _sentence_spans(text: str):
    start = 0
    for boundary in _SENTENCE_BOUNDARY_RE.finditer(text):
        if boundary.start() > start:
            yield start, boundary.start()
        start = boundary.end()
    if start < len(text):
        yield start, len(text)


def _generation_scopes(text: str) -> list[tuple[int, int]]:
    """Return portions governed by a non-negated production verb."""

    scopes: list[tuple[int, int]] = []
    for sentence_start, sentence_end in _sentence_spans(text):
        sentence = text[sentence_start:sentence_end]
        actions = list(_ACTION_RE.finditer(sentence))
        for index, action in enumerate(actions):
            comma_start = 0
            preceding_commas = list(_COMMA_RE.finditer(sentence, 0, action.start()))
            if preceding_commas:
                comma_start = preceding_commas[-1].end()
            local_end_match = _COMMA_RE.search(sentence, action.end())
            local_end = local_end_match.start() if local_end_match else len(sentence)
            local_clause = sentence[comma_start:local_end]
            if _NEGATED_PRODUCTION_RE.search(local_clause):
                continue

            previous_end = actions[index - 1].end() if index else 0
            scope_start = max(previous_end, comma_start)
            # Formats mentioned by a preceding read/analyse clause describe
            # inputs, not outputs ("读取 Excel 并生成 PDF").
            if _INPUT_CONTEXT_RE.search(sentence[scope_start:action.start()]):
                scope_start = action.start()
            scope_end = actions[index + 1].start() if index + 1 < len(actions) else len(sentence)
            after_action = sentence[action.end():scope_end]
            boundary = _OUTPUT_SCOPE_BOUNDARY_RE.search(after_action)
            if boundary:
                scope_end = action.end() + boundary.start()
            scope = sentence[scope_start:scope_end]
            if action.lastgroup == "weak":
                if not _FILE_NOUN_RE.search(scope) and not _COUNT_TARGET_RE.search(scope):
                    continue
            scopes.append((sentence_start + scope_start, sentence_start + scope_end))
    return scopes


def _format_mentions(text: str) -> list[tuple[int, str]]:
    mentions: list[tuple[int, str]] = []
    for start, end in _generation_scopes(text):
        scope = text[start:end]
        input_regions = list(_FORMAT_INPUT_REGION_RE.finditer(scope))
        for file_format, pattern in _FORMAT_PATTERNS.items():
            for match in pattern.finditer(scope):
                # "生成基于 Excel 数据的 PDF 报告" has two format
                # mentions, but only PDF is an output.
                if any(
                    region.start() <= match.start() < region.end()
                    for region in input_regions
                ):
                    continue
                mentions.append((start + match.start(), file_format))
    mentions.sort(key=lambda item: item[0])
    return mentions


def _truncate_utf8(value: str, max_bytes: int) -> str:
    encoded = value.encode("utf-8")
    if len(encoded) <= max_bytes:
        return value
    return encoded[:max_bytes].decode("utf-8", errors="ignore")


def _safe_filename(value: str | None, file_format: str) -> str:
    spec = FORMAT_SPECS[file_format]
    raw = unicodedata.normalize("NFKC", str(value or "")).replace("\\", "/")
    raw = Path(raw).name
    raw = "".join(char for char in raw if char >= " " and char != "\x7f")
    raw = re.sub(r'[<>:"/\\|?*]+', "-", raw)
    raw = re.sub(r"\s+", " ", raw).strip(" .-")
    if raw:
        suffix = Path(raw).suffix.lower().lstrip(".")
        stem = Path(raw).stem if suffix else raw
    else:
        stem = DEFAULT_FILE_STEM
    stem = _truncate_utf8(stem.strip(" .-") or DEFAULT_FILE_STEM, 140)
    return f"{stem}.{spec.extension}"


def _requested_filename(text: str) -> str | None:
    match = _FILENAME_RE.search(text)
    if not match:
        return None
    return (match.group("quoted") or match.group("bare") or "").strip()


def detect_file_artifact_requests(message: str | None) -> list[ArtifactRequest]:
    """Detect explicit, non-negated file creation requests.

    Reading a file, asking why creation is unavailable, or asking to resend a
    previously generated file does not contain an active production scope and
    therefore returns an empty list.
    """

    text = str(message or "")
    raw_filename = _requested_filename(text)
    ordered_formats: list[str] = []
    seen: set[str] = set()
    for _, file_format in _format_mentions(text):
        if file_format in seen:
            continue
        seen.add(file_format)
        ordered_formats.append(file_format)
    if len(ordered_formats) > MAX_FILE_ARTIFACTS:
        raise ValueError(f"一次最多生成 {MAX_FILE_ARTIFACTS} 个文件")
    return [
        ArtifactRequest(
            format=file_format,
            filename=_safe_filename(raw_filename, file_format),
        )
        for file_format in ordered_formats
    ]


def _clean_explicit_content(value: str) -> str:
    value = str(value or "").strip()
    quote_pairs = {"\"": "\"", "'": "'", "“": "”", "‘": "’", "`": "`"}
    if value[:1] in quote_pairs:
        closing = quote_pairs[value[0]]
        closing_index = value.find(closing, 1)
        if closing_index >= 1:
            return value[1:closing_index][:MAX_ARTIFACT_TEXT_CHARS]
    value = _FOLLOW_UP_RE.split(value, maxsplit=1)[0]
    value = re.split(
        r"\s*(?:就好|即可|就可以|就行|就够了)(?:\s|[。！!？?，,、]|$).*",
        value,
        maxsplit=1,
    )[0]
    return value.strip().strip("\"'“”‘’` ")[:MAX_ARTIFACT_TEXT_CHARS]


def extract_explicit_file_content(message: str | None) -> str | None:
    """Extract user-supplied literal/fenced content without consulting a model."""

    text = str(message or "")
    fenced = _FENCED_BLOCK_RE.search(text)
    if fenced:
        prefix = text[:fenced.start()]
        explicit_intro = re.search(
            r"(?:content|正文|数据|代码)(?:里面|里|中)?\s*"
            r"(?:写(?:入|上)?|填写|输入|为|是|包含|放入?|如下|：|:)\s*$"
            r"|(?:文件|文档|页面|脚本)\s*(?:如下|：|:)\s*$"
            r"|\b(?:content|text|body|data|code)\s*"
            r"(?:is|be|below|follows|:|=)\s*$",
            prefix[-240:],
            re.IGNORECASE,
        )
        language = (fenced.group("language") or "").strip().casefold()
        requested_formats = {
            file_format for _, file_format in _format_mentions(text)
        }
        language_matches_output = any(
            language in _FENCE_LANGUAGE_FORMATS.get(file_format, set())
            for file_format in requested_formats
        )
        # A fenced input being analysed is not the requested output body.
        # Accept it deterministically only when explicitly introduced as
        # content or when its language matches a requested output format.
        if explicit_intro or (language and language_matches_output):
            body = fenced.group("body")[:MAX_ARTIFACT_TEXT_CHARS]
            return body if body.strip() else None
    for pattern in _CONTENT_PATTERNS:
        match = pattern.search(text)
        if match:
            raw_value = match.group("value")
            stripped = str(raw_value or "").strip()
            quote_pairs = {"\"": "\"", "'": "'", "“": "”", "‘": "’", "`": "`"}
            explicitly_quoted = bool(
                stripped[:1] in quote_pairs
                and quote_pairs[stripped[0]] in stripped[1:]
            )
            value = _clean_explicit_content(raw_value)
            derived_instruction = re.search(
                r"(?:"
                r"(?:根据|基于|参考|读取|阅读|查看|解析|分析|总结|提取|整理|改写|翻译)"
                r".{0,24}(?:上传|附件|文件|文档|表格|Excel|Word|PDF|PPT|数据|对话)"
                r"|(?:上传|附件|文件|文档|表格|Excel|Word|PDF|PPT|数据|对话)"
                r".{0,24}(?:总结|分析|读取|解析|提取|整理|改写|翻译)"
                r"|\b(?:based\s+on|summari[sz]e|analy[sz]e|extract|rewrite|translate)"
                r".{0,32}\b(?:attached|uploaded|file|document|spreadsheet|data|chat)\b"
                r")",
                value,
                re.IGNORECASE,
            )
            if derived_instruction and not explicitly_quoted:
                return None
            if value:
                return value
    # Preserve deterministic existing-workbook cell updates such as
    # "在上传的 Excel 的 B2 单元格写入 hello 并导出".
    excel_value = extract_excel_a1_content(text)
    if excel_value:
        return excel_value
    return None


def _fenced_content_for_format(value: str, file_format: str) -> str | None:
    candidates = list(_FENCED_BLOCK_RE.finditer(str(value or "")))
    if not candidates:
        return None
    languages = _FENCE_LANGUAGE_FORMATS.get(file_format, set())
    for match in candidates:
        language = (match.group("language") or "").strip().casefold()
        if language in languages:
            return match.group("body")
    if len(candidates) == 1:
        return candidates[0].group("body")
    return None


def _content_for_format(
    *,
    file_format: str,
    explicit_content: str | None,
    model_reply: str | None,
) -> str:
    source = explicit_content if explicit_content is not None else str(model_reply or "").strip()
    if not source:
        return ""
    fenced = _fenced_content_for_format(source, file_format)
    value = fenced if fenced is not None else source
    value = value.replace("\x00", "")
    return value[:MAX_ARTIFACT_TEXT_CHARS]


def _text_bytes(content: str) -> bytes:
    return (content.rstrip() + "\n").encode("utf-8")


def _render_json(content: str) -> bytes:
    try:
        payload = json.loads(content)
    except (json.JSONDecodeError, TypeError):
        payload = {"content": content}
    return (
        json.dumps(payload, ensure_ascii=False, indent=2, allow_nan=False) + "\n"
    ).encode("utf-8")


def _safe_spreadsheet_cell(value: object) -> str:
    text = str(value if value is not None else "").replace("\x00", "")
    leading = text.lstrip()
    if leading.startswith(("=", "+", "-", "@")):
        prefix_len = len(text) - len(leading)
        text = f"{text[:prefix_len]}'{leading}"
    return text


def _render_csv(content: str) -> bytes:
    if not any(separator in content for separator in (",", "\t", ";", "\n", "\r")):
        rows = [[content]]
    else:
        sample = content[:4096]
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t")
        except csv.Error:
            dialect = csv.excel
        rows = list(csv.reader(io.StringIO(content), dialect))
        if not rows:
            rows = [[""]]
    stream = io.StringIO(newline="")
    writer = csv.writer(stream, dialect=csv.excel, lineterminator="\n")
    for row in rows:
        writer.writerow([_safe_spreadsheet_cell(cell) for cell in row])
    return stream.getvalue().encode("utf-8-sig")


_STATIC_HTML_TAGS = {
    "a",
    "abbr",
    "article",
    "aside",
    "b",
    "blockquote",
    "br",
    "caption",
    "code",
    "col",
    "colgroup",
    "dd",
    "del",
    "details",
    "div",
    "dl",
    "dt",
    "em",
    "figcaption",
    "figure",
    "footer",
    "h1",
    "h2",
    "h3",
    "h4",
    "h5",
    "h6",
    "header",
    "hr",
    "i",
    "img",
    "ins",
    "kbd",
    "li",
    "main",
    "mark",
    "nav",
    "ol",
    "p",
    "pre",
    "s",
    "samp",
    "section",
    "small",
    "span",
    "strong",
    "sub",
    "summary",
    "sup",
    "table",
    "tbody",
    "td",
    "tfoot",
    "th",
    "thead",
    "time",
    "tr",
    "u",
    "ul",
    "var",
}
_STATIC_HTML_GLOBAL_ATTRIBUTES = {
    "class",
    "dir",
    "id",
    "lang",
    "role",
    "style",
    "title",
}
_STATIC_HTML_TAG_ATTRIBUTES = {
    "a": {"href"},
    "img": {"alt", "height", "src", "width"},
    "ol": {"reversed", "start", "type"},
    "td": {"colspan", "rowspan"},
    "th": {"colspan", "rowspan", "scope"},
    "time": {"datetime"},
}
_STATIC_HTML_CSS_PROPERTIES = {
    "align-items",
    "background-color",
    "border",
    "border-bottom",
    "border-collapse",
    "border-color",
    "border-left",
    "border-radius",
    "border-right",
    "border-style",
    "border-top",
    "border-width",
    "box-sizing",
    "color",
    "display",
    "flex",
    "flex-direction",
    "flex-wrap",
    "font-family",
    "font-size",
    "font-style",
    "font-weight",
    "gap",
    "grid-template-columns",
    "height",
    "justify-content",
    "line-height",
    "list-style",
    "margin",
    "margin-bottom",
    "margin-left",
    "margin-right",
    "margin-top",
    "max-width",
    "min-width",
    "overflow-wrap",
    "padding",
    "padding-bottom",
    "padding-left",
    "padding-right",
    "padding-top",
    "text-align",
    "text-decoration",
    "vertical-align",
    "white-space",
    "width",
    "word-break",
}
_ACTIVE_HTML_BLOCK_RE = re.compile(
    r"<\s*(script|style|noscript|template|iframe|object|embed|applet|svg|math)"
    r"\b[^>]*>.*?<\s*/\s*\1\s*>",
    re.IGNORECASE | re.DOTALL,
)
_STATIC_HTML_CSP = (
    "default-src 'none'; img-src data:; style-src 'unsafe-inline'; "
    "font-src data:; form-action 'none'; base-uri 'none'; "
    "frame-src 'none'; object-src 'none'"
)


def _static_html_attribute_filter(tag: str, name: str, value: str) -> bool:
    lowered_name = str(name or "").casefold()
    if lowered_name.startswith("on"):
        return False
    allowed = (
        lowered_name in _STATIC_HTML_GLOBAL_ATTRIBUTES
        or lowered_name in _STATIC_HTML_TAG_ATTRIBUTES.get(tag, set())
        or lowered_name.startswith("aria-")
    )
    if not allowed:
        return False
    if tag == "img" and lowered_name == "src":
        normalized = html.unescape(value or "").strip().casefold()
        return bool(
            re.match(
                r"^data:image/(?:png|jpeg|gif|webp);base64,[a-z0-9+/=\s]+$",
                normalized,
            )
        )
    if tag == "a" and lowered_name == "href":
        normalized = html.unescape(value or "").strip().casefold()
        return (
            normalized.startswith("#")
            or normalized.startswith("https://")
            or normalized.startswith("mailto:")
        )
    return True


def _sanitize_static_html(value: str) -> str:
    source = _ACTIVE_HTML_BLOCK_RE.sub("", value)
    try:
        import bleach
        from bleach.css_sanitizer import CSSSanitizer

        return bleach.clean(
            source,
            tags=_STATIC_HTML_TAGS,
            attributes=_static_html_attribute_filter,
            protocols={"https", "mailto", "data"},
            css_sanitizer=CSSSanitizer(
                allowed_css_properties=_STATIC_HTML_CSS_PROPERTIES,
            ),
            strip=True,
            strip_comments=True,
        )
    except ImportError:
        # Missing optional runtime dependency must fail closed: preserve the
        # source as visible text instead of serving active markup.
        return f"<pre>{html.escape(source)}</pre>"


def _render_html(content: str) -> bytes:
    if re.search(
        r"<\s*(?:!doctype|html|head|body|main|article|section|div|p|h[1-6])\b",
        content,
        re.IGNORECASE,
    ):
        body = _sanitize_static_html(content)
    else:
        body = f"<pre>{html.escape(content)}</pre>"
    document = (
        "<!doctype html>\n"
        '<html lang="zh-CN">\n'
        "<head>\n"
        '<meta charset="utf-8">\n'
        f'<meta http-equiv="Content-Security-Policy" content="{_STATIC_HTML_CSP}">\n'
        '<meta name="viewport" content="width=device-width, initial-scale=1">\n'
        "<title>小策bot 产出</title>\n"
        "</head>\n"
        f"<body>{body}</body>\n"
        "</html>\n"
    )
    return document.encode("utf-8")


def _render_xml(content: str) -> bytes:
    candidate = content.strip()
    if candidate.startswith("<"):
        try:
            from defusedxml import ElementTree as SafeElementTree

            SafeElementTree.fromstring(candidate)
            document = candidate
        except Exception:
            document = f"<document><content>{html.escape(content)}</content></document>"
    else:
        document = f"<document><content>{html.escape(content)}</content></document>"
    return (document.rstrip() + "\n").encode("utf-8")


def _render_docx(content: str) -> bytes:
    from docx import Document

    document = Document()
    emitted = False
    for raw_line in content.splitlines() or [content]:
        line = raw_line.rstrip()
        heading = re.match(r"^(#{1,6})\s+(.+)$", line)
        bullet = re.match(r"^\s*[-*+]\s+(.+)$", line)
        numbered = re.match(r"^\s*\d+[.)]\s+(.+)$", line)
        if heading:
            document.add_heading(heading.group(2), level=min(len(heading.group(1)), 6))
            emitted = True
        elif bullet:
            document.add_paragraph(bullet.group(1), style="List Bullet")
            emitted = True
        elif numbered:
            document.add_paragraph(numbered.group(1), style="List Number")
            emitted = True
        elif line:
            document.add_paragraph(line)
            emitted = True
        elif emitted:
            document.add_paragraph("")
    if not emitted:
        document.add_paragraph("")
    stream = io.BytesIO()
    document.save(stream)
    return stream.getvalue()


def _presentation_sections(content: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    title = "小策bot 产出"
    lines: list[str] = []
    for raw_line in content.splitlines():
        heading = re.match(r"^#{1,2}\s+(.+)$", raw_line.strip())
        if heading:
            if lines or sections:
                sections.append((title, lines))
            title = heading.group(1).strip() or title
            lines = []
        elif raw_line.strip() == "---" and lines:
            sections.append((title, lines))
            title = "续"
            lines = []
        elif raw_line.strip():
            lines.append(re.sub(r"^\s*(?:[-*+]|\d+[.)])\s+", "", raw_line.strip()))
    if lines or not sections:
        sections.append((title, lines or [content.strip()]))
    return sections[:50]


def _render_pptx(content: str) -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    for title, lines in _presentation_sections(content):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, line in enumerate(lines or [""]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _pdf_font_name(content: str) -> str:
    if all(ord(char) < 256 for char in content):
        return "Helvetica"
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont

    font_name = "STSong-Light"
    try:
        pdfmetrics.getFont(font_name)
    except KeyError:
        pdfmetrics.registerFont(UnicodeCIDFont(font_name))
    return font_name


def _wrap_pdf_line(line: str, width: int) -> list[str]:
    if not line:
        return [""]
    return [line[index:index + width] for index in range(0, len(line), width)]


def _render_pdf(content: str) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    stream = io.BytesIO()
    document = canvas.Canvas(stream, pagesize=A4, pageCompression=1)
    document.setTitle("小策bot 产出")
    font_name = _pdf_font_name(content)
    _, page_height = A4
    x = 48
    y = page_height - 52
    line_height = 16
    document.setFont(font_name, 10.5)
    wrap_width = 46 if font_name != "Helvetica" else 92
    for source_line in content.splitlines() or [content]:
        for line in _wrap_pdf_line(source_line, wrap_width):
            if y < 52:
                document.showPage()
                document.setFont(font_name, 10.5)
                y = page_height - 52
            document.drawString(x, y, line)
            y -= line_height
    document.save()
    return stream.getvalue()


_RENDERERS: dict[str, Callable[[str], bytes]] = {
    "md": _text_bytes,
    "txt": _text_bytes,
    "html": _render_html,
    "json": _render_json,
    "csv": _render_csv,
    "docx": _render_docx,
    "pdf": _render_pdf,
    "pptx": _render_pptx,
    "py": _text_bytes,
    "js": _text_bytes,
    "ts": _text_bytes,
    "jsx": _text_bytes,
    "tsx": _text_bytes,
    "css": _text_bytes,
    "sql": _text_bytes,
    "yaml": _text_bytes,
    "xml": _render_xml,
}


def _store_artifact(*, user_id: int, request: ArtifactRequest, data: bytes) -> dict:
    if not data:
        raise ValueError(f"{request.format} 文件内容为空")
    if len(data) > MAX_GENERATED_FILE_BYTES:
        raise ValueError(
            f"{request.filename} 超过 "
            f"{MAX_GENERATED_FILE_BYTES // (1024 * 1024)}MB 生成上限",
        )

    root = attachments_root(user_id)
    root.mkdir(parents=True, exist_ok=True)
    stored_id = f"{uuid.uuid4().hex}_{request.filename}"
    target = root / stored_id
    temporary_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(
            mode="wb",
            prefix=".xiaoce-file-",
            suffix=f".{request.format}.tmp",
            dir=root,
            delete=False,
        ) as temporary:
            temporary.write(data)
            temporary.flush()
            os.fsync(temporary.fileno())
            temporary_path = Path(temporary.name)
        os.replace(temporary_path, target)
        temporary_path = None
    finally:
        if temporary_path is not None:
            temporary_path.unlink(missing_ok=True)

    spec = FORMAT_SPECS[request.format]
    return {
        "id": stored_id,
        "name": request.filename,
        "size": target.stat().st_size,
        "mime": spec.mime,
        "has_text": spec.text_file,
        "is_image": False,
        "is_file": True,
    }


def _rename_delegated_artifact(
    *,
    user_id: int,
    item: dict,
    filename: str,
) -> dict:
    current_id = str(item.get("id") or item.get("stored_id") or "")
    current_path = resolve_attachment_path(user_id, current_id)
    if current_path is None:
        raise ValueError("Excel 生成结果文件不可用")
    if current_path.stat().st_size > MAX_GENERATED_FILE_BYTES:
        current_path.unlink(missing_ok=True)
        raise ValueError("Excel 生成结果超过文件大小上限")
    if str(item.get("name") or "") == filename:
        return dict(item)

    new_id = f"{uuid.uuid4().hex}_{filename}"
    new_path = attachments_root(user_id) / new_id
    os.replace(current_path, new_path)
    renamed = {
        **item,
        "id": new_id,
        "name": filename,
        "size": new_path.stat().st_size,
    }
    if "stored_id" in item:
        renamed["stored_id"] = new_id
    return renamed


def _cleanup_generated(items: list[dict], user_id: int) -> None:
    for item in items:
        stored_id = str(item.get("id") or item.get("stored_id") or "")
        path = resolve_attachment_path(user_id, stored_id)
        if path is not None:
            path.unlink(missing_ok=True)


def maybe_generate_file_artifacts(
    *,
    user_id: int,
    request_text: str | None,
    model_reply: str | None,
    source_attachments: list[dict] | None = None,
) -> list[dict]:
    """Create all explicitly requested files, at most five per turn.

    XLSX remains delegated to the hardened Excel implementation so existing
    workbook-copy/cell-edit behaviour is retained.
    """

    requests = detect_file_artifact_requests(request_text)
    if not requests:
        return []
    explicit_content = extract_explicit_file_content(request_text)
    if explicit_content is None and not str(model_reply or "").strip():
        return []

    generated: list[dict] = []
    try:
        for request in requests:
            if request.format == "xlsx":
                excel_items = maybe_generate_excel_artifact(
                    user_id=user_id,
                    request_text=request_text,
                    model_reply=model_reply,
                    source_attachments=source_attachments,
                )
                for item in excel_items:
                    try:
                        renamed = _rename_delegated_artifact(
                            user_id=user_id,
                            item=item,
                            filename=request.filename,
                        )
                    except Exception:
                        _cleanup_generated([item], user_id)
                        raise
                    generated.append(renamed)
                continue

            content = _content_for_format(
                file_format=request.format,
                explicit_content=explicit_content,
                model_reply=model_reply,
            )
            if not content:
                raise ValueError(f"{request.format} 文件缺少可生成内容")
            renderer = _RENDERERS[request.format]
            generated.append(
                _store_artifact(
                    user_id=user_id,
                    request=request,
                    data=renderer(content),
                ),
            )
        return generated
    except Exception:
        _cleanup_generated(generated, user_id)
        raise


__all__ = [
    "ArtifactRequest",
    "DEFAULT_FILE_STEM",
    "FORMAT_SPECS",
    "MAX_FILE_ARTIFACTS",
    "detect_file_artifact_requests",
    "extract_explicit_file_content",
    "maybe_generate_file_artifacts",
]
